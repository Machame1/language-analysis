import os
import re
from flask import Flask, render_template, request, redirect, url_for
import PyPDF2
import docx
from textblob import TextBlob
import nltk
from deep_translator import GoogleTranslator  # Import deep-translator

from pptx import Presentation
from collections import defaultdict
from nltk.corpus import stopwords

nltk.download('punkt')
nltk.download('stopwords')

app = Flask(__name__)

import nltk
from nltk.corpus import stopwords
from collections import Counter


def translate_text(text, target_lang):
    """Translate text using deep-translator."""
    try:
        translated_text = GoogleTranslator(source='auto', target=target_lang).translate(text)
        return translated_text
    except Exception as e:
        return str(e)

def generate_title(text):
    # Tokenize the text into words
    words = nltk.word_tokenize(text.lower())
    
    # Filter out stop words and non-alphanumeric words
    stop_words = set(stopwords.words('english'))
    filtered_words = [word for word in words if word.isalnum() and word not in stop_words]

    # Count word frequencies
    word_counts = Counter(filtered_words)
    
    # Get the most common words
    most_common = word_counts.most_common(3)
    
    # Generate a title from the most common words
    if most_common:
        title = ' '.join([word for word, _ in most_common]).title()
    else:
        title = "Untitled Document"

    return title



def read_text_from_pdf(file_path):
    text = ''
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() or '' 
    return text

def read_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return '\n'.join([para.text for para in doc.paragraphs])

def read_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def read_text_from_pptx(file_path):
    text = ''
    presentation = Presentation(file_path)
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = f"Slide {slide_number}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        text += slide_text + "\n"
    return text

def remove_noise(text):
    return re.sub(r'[^\w\s]', '', text.lower())

def overall_sentiment(text):
    cleaned_text = remove_noise(text)
    analysis = TextBlob(cleaned_text)
    polarity = analysis.sentiment.polarity
    return 'Positive' if polarity > 0 else 'Negative' if polarity < 0 else 'Neutral'

def summarize_overall_content(text, num_sentences=3):
    # Split the text into sentences
    sentences = nltk.sent_tokenize(text)

    # Create a frequency distribution of words
    word_frequencies = defaultdict(int)
    stop_words = set(stopwords.words('english'))

    for word in nltk.word_tokenize(text.lower()):
        if word.isalnum() and word not in stop_words:
            word_frequencies[word] += 1

    # Score sentences based on the frequency of words
    sentence_scores = defaultdict(int)
    for sentence in sentences:
        for word in nltk.word_tokenize(sentence.lower()):
            if word in word_frequencies:
                sentence_scores[sentence] += word_frequencies[word]

    # Select the top N sentences to form the summary
    summarized_sentences = sorted(sentence_scores, key=sentence_scores.get, reverse=True)[:num_sentences]

    # Create a summary by combining the selected sentences
    summary = ' '.join(summarized_sentences)
    
    return summary


@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return redirect(request.url)
    
    file = request.files['file']
    
    if file:
        os.makedirs('uploads', exist_ok=True)
        file_path = os.path.join('uploads', file.filename)
        file.save(file_path)
        return redirect(url_for('upload_success', filename=file.filename))
    
    return "No file uploaded", 400

@app.route('/upload_success/<filename>')
def upload_success(filename):
    file_path = os.path.join('uploads', filename)
    text = read_text_from_pdf(file_path) if filename.endswith('.pdf') else \
           read_text_from_docx(file_path) if filename.endswith('.docx') else \
           read_text_from_txt(file_path) if filename.endswith('.txt') else \
           read_text_from_pptx(file_path)

    initial_content = text[:300]
    remaining_content = text[300:] if len(text) > 300 else ""

    return render_template('success_view.html', initial_content=initial_content, remaining_content=remaining_content, filename=filename)

@app.route('/analyze/<filename>')
def analyze_sentiment_route(filename):
    file_path = os.path.join('uploads', filename)
    text = read_text_from_pdf(file_path) if filename.endswith('.pdf') else \
           read_text_from_docx(file_path) if filename.endswith('.docx') else \
           read_text_from_txt(file_path) if filename.endswith('.txt') else \
           read_text_from_pptx(file_path)

    sentiment_label = overall_sentiment(text)

    initial_content = text[:300]
    remaining_content = text[300:] if len(text) > 300 else ""
    
    return render_template('result.html', results=[("Overall Sentiment", sentiment_label)],
                           filename=filename,
                          
                           initial_content=initial_content,
                           remaining_content=remaining_content)

@app.route('/remove_noise/<filename>')
def remove_noise_route(filename):
    file_path = os.path.join('uploads', filename)
    text = read_text_from_pdf(file_path) if filename.endswith('.pdf') else \
           read_text_from_docx(file_path) if filename.endswith('.docx') else \
           read_text_from_txt(file_path) if filename.endswith('.txt') else \
           read_text_from_pptx(file_path)

    cleaned_text = remove_noise(text)
    sentiment_label = overall_sentiment(cleaned_text)
    
    initial_content = text[:300]
    remaining_content = text[300:] if len(text) > 300 else ""
    
    return render_template('result.html', 
                           results=[("Cleaned Text Sentiment", sentiment_label)],
                           filename=filename,   
                           cleaned_content=cleaned_text
                           )

@app.route('/summary/<filename>', methods=['GET', 'POST'])
def summarize_route(filename):
    file_path = os.path.join('uploads', filename)
    text = read_text_from_pdf(file_path) if filename.endswith('.pdf') else \
           read_text_from_docx(file_path) if filename.endswith('.docx') else \
           read_text_from_txt(file_path) if filename.endswith('.txt') else \
           read_text_from_pptx(file_path)

    title = generate_title(text)
    summary = summarize_overall_content(text, num_sentences=5)

    translated_summary = None  # Initialize translated summary variable

    if request.method == 'POST':
        target_lang = request.form['language']
        translated_summary = translate_text(summary, target_lang)  # Translate on form submission

    return render_template('summary.html', title=title, summary=summary, translated_summary=translated_summary)



if __name__ == '__main__':
    os.makedirs('uploads', exist_ok=True)
    app.run(debug=True, port=4000)
